#!/usr/bin/env python3
"""Convert a PowerPoint (.pptx) file to nested HTML unordered lists.

Usage:  python3 slides_to_html.py presentation.pptx
Output: HTML printed to stdout — redirect to a file as needed.

Requires: pip3 install python-pptx
"""

import sys
import html as html_lib
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER


def esc(text):
    return html_lib.escape(str(text).strip())


def is_title(shape):
    """Return True if this shape is a slide title placeholder."""
    try:
        return shape.placeholder_format.type in (
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
        )
    except (AttributeError, TypeError, ValueError):
        return False


def get_slide_title(slide):
    """Return the title text of a slide."""
    for shape in slide.shapes:
        if is_title(shape) and shape.has_text_frame:
            return esc(shape.text_frame.text)
    # Fallback: first non-empty text shape
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                return esc(text)
    return "Untitled Slide"


def build_list(items, depth=0):
    """Recursively convert (level, text) tuples into nested <ul><li> HTML.

    PowerPoint paragraph levels (0, 1, 2 ...) map directly to nesting depth.
    Items at a deeper level than the current base level become children of
    the preceding item.
    """
    if not items:
        return ""

    pad = "    " * depth
    out = [f"{pad}<ul>"]
    i = 0
    base = items[0][0]

    while i < len(items):
        lvl, text = items[i]
        if lvl != base:
            i += 1
            continue

        # Collect everything between this item and the next same-level item as children
        j = i + 1
        while j < len(items) and items[j][0] > base:
            j += 1

        children = items[i + 1 : j]
        if children:
            out.append(f"{pad}    <li>{text}")
            out.append(build_list(children, depth + 1))
            out.append(f"{pad}    </li>")
        else:
            out.append(f"{pad}    <li>{text}</li>")
        i = j

    out.append(f"{pad}</ul>")
    return "\n".join(out)


def table_to_html(table, depth=0):
    """Convert a pptx table to nested <ul><li> elements.

    The first row provides the column headers, each becoming a parent <li>.
    All values in that column become a nested <ul> under their header.
    """
    if not table.rows:
        return ""

    pad = "    " * depth
    rows = list(table.rows)
    headers = [esc(cell.text) for cell in rows[0].cells]
    data_rows = [[esc(cell.text) for cell in row.cells] for row in rows[1:]]

    out = [f"{pad}<ul>"]
    for col_idx, header in enumerate(headers):
        if not header:
            continue
        col_values = [
            row[col_idx] for row in data_rows
            if col_idx < len(row) and row[col_idx]
        ]
        if col_values:
            out.append(f"{pad}    <li>{header}")
            out.append(f"{pad}        <ul>")
            for val in col_values:
                out.append(f"{pad}            <li>{val}</li>")
            out.append(f"{pad}        </ul>")
            out.append(f"{pad}    </li>")
        else:
            out.append(f"{pad}    <li>{header}</li>")
    out.append(f"{pad}</ul>")
    return "\n".join(out)


def slide_content_to_html(slide, depth=0):
    """Extract all non-title content from a slide and return HTML."""
    parts = []

    for shape in slide.shapes:
        if is_title(shape):
            continue

        if shape.has_text_frame:
            items = []
            for para in shape.text_frame.paragraphs:
                text = esc("".join(run.text for run in para.runs))
                if text:
                    items.append((para.level, text))
            if items:
                parts.append(build_list(items, depth))

        # Tables: check with hasattr since only GraphicFrame shapes have has_table
        elif hasattr(shape, "has_table") and shape.has_table:
            parts.append(table_to_html(shape.table, depth))

    return "\n".join(parts)


def convert(filepath):
    prs = Presentation(filepath)
    print("<ul>")
    for slide in prs.slides:
        title = get_slide_title(slide)
        content = slide_content_to_html(slide, depth=2)
        if content.strip():
            print(f"    <li>{title}")
            print(content)
            print(f"    </li>")
        else:
            print(f"    <li>{title}</li>")
    print("</ul>")


if __name__ == "__main__":
    if len(sys.argv) != 2:
        sys.exit(f"Usage: python3 {sys.argv[0]} <file.pptx>")
    convert(sys.argv[1])
